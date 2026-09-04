#!/usr/bin/env python3
"""
Build a LLaMAR-style teaser deck for the JetRover project (English, for interviews).
Re-run any time:  python3 build_deck.py
Output: JetRover_teaser.pptx  in the same folder.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------- design tokens ----------
BAR     = RGBColor(0x8C, 0x1D, 0x2A)   # maroon bottom bar
TITLE   = RGBColor(0x1A, 0x1A, 0x1A)   # near-black title
INK     = RGBColor(0x22, 0x22, 0x22)
MUTE    = RGBColor(0x66, 0x66, 0x66)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BORDER  = RGBColor(0x5A, 0x5A, 0x5A)

C_BLUE   = RGBColor(0xDA, 0xE8, 0xFC)
C_PURPLE = RGBColor(0xE6, 0xDA, 0xF2)
C_GRAY   = RGBColor(0xED, 0xED, 0xED)
C_GREEN  = RGBColor(0xD5, 0xE8, 0xD4)
C_YELLOW = RGBColor(0xFF, 0xF2, 0xCC)
C_ORANGE = RGBColor(0xFF, 0xE5, 0x99)
C_TAN    = RGBColor(0xFC, 0xE5, 0xCD)
C_PINK   = RGBColor(0xF5, 0xC6, 0xD0)
C_PH     = RGBColor(0xF3, 0xF3, 0xF3)   # image placeholder

FONT = "Arial"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
PGW, PGH = 13.333, 7.5


def _no_shadow(shp):
    try:
        shp.shadow.inherit = False
    except Exception:
        pass


def slide():
    return prs.slides.add_slide(BLANK)


def footer(s, n):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(PGH - 0.32),
                             Inches(PGW), Inches(0.32))
    bar.fill.solid(); bar.fill.fore_color.rgb = BAR
    bar.line.fill.background(); _no_shadow(bar)
    # wordmark
    wm = s.shapes.add_textbox(Inches(0.32), Inches(PGH - 0.34), Inches(6), Inches(0.32))
    p = wm.text_frame.paragraphs[0]; r = p.add_run(); r.text = "JetRover"
    r.font.name = FONT; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE
    # page number
    pg = s.shapes.add_textbox(Inches(PGW - 1.1), Inches(PGH - 0.34), Inches(0.8), Inches(0.32))
    p = pg.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = str(n)
    r.font.name = FONT; r.font.size = Pt(12); r.font.color.rgb = WHITE


def title(s, text, size=40):
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = TITLE
    return tb


def box(s, x, y, w, h, text, fill, fs=15, bold=False, rounded=True,
        fcolor=INK, line=BORDER, lw=1.0, align=PP_ALIGN.CENTER):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    _no_shadow(shp)
    tf = shp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = FONT; r.font.size = Pt(fs); r.font.bold = bold; r.font.color.rgb = fcolor
    return shp


def note(s, x, y, w, h, text, fs=12):
    """thin no-fill annotation box (like LLaMAR side notes)"""
    return box(s, x, y, w, h, text, WHITE, fs=fs, bold=False, rounded=False,
               fcolor=INK, line=BORDER, lw=0.75)


def placeholder(s, x, y, w, h, label):
    shp = box(s, x, y, w, h, label, C_PH, fs=14, bold=False, rounded=True,
              fcolor=MUTE, line=RGBColor(0xBB, 0xBB, 0xBB), lw=1.0)
    return shp


def arrow(s, x1, y1, x2, y2, color=RGBColor(0x55, 0x55, 0x55), width=1.5):
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                  Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    _no_shadow(conn)
    return conn


def alabel(s, x, y, text, fs=11):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(1.6), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(fs); r.font.color.rgb = MUTE


def bullets(s, x, y, w, items, fs=22, gap=10):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            txt, sub = it
            r = p.add_run(); r.text = "•  " + txt
            r.font.name = FONT; r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = INK
            p2 = tf.add_paragraph(); p2.space_after = Pt(gap)
            r2 = p2.add_run(); r2.text = "     " + sub
            r2.font.name = FONT; r2.font.size = Pt(fs - 5); r2.font.color.rgb = MUTE
        else:
            r = p.add_run(); r.text = "•  " + it
            r.font.name = FONT; r.font.size = Pt(fs); r.font.color.rgb = INK
    return tb


def big_text(s, x, y, w, text, fs=30, bold=True, color=INK, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(fs); r.font.bold = bold; r.font.color.rgb = color
    return tb


# ============================================================ SLIDE 1 — Title
s = slide()
tb = s.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(12), Inches(2.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "JetRover"
r.font.name = FONT; r.font.size = Pt(60); r.font.bold = True; r.font.color.rgb = TITLE
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r = p2.add_run(); r.text = "A Language-Guided Mobile Manipulator"
r.font.name = FONT; r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = BAR
sub = s.shapes.add_textbox(Inches(1.5), Inches(4.15), Inches(10.3), Inches(1.0))
tf = sub.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Closing the perception → planning → action → memory loop on real hardware"
r.font.name = FONT; r.font.size = Pt(19); r.font.color.rgb = MUTE
au = s.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.8))
p = au.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "[ Your Name ]   ·   2026"
r.font.name = FONT; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = INK
pf = s.shapes.add_textbox(Inches(1.5), Inches(5.85), Inches(10.3), Inches(0.6))
p = pf.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Hiwonder JetRover  ·  mecanum base + 5-DOF arm  ·  Orbbec RGB-D camera  ·  ROS"
r.font.name = FONT; r.font.size = Pt(13); r.font.color.rgb = MUTE
footer(s, 1)

# ============================================================ SLIDE 2 — Motivation
s = slide()
title(s, "Motivation")
big_text(s, 1.2, 1.9, 11, "“I want a coke — can you get me one?”",
         fs=34, bold=True, color=INK, align=PP_ALIGN.CENTER)
big_text(s, 1.6, 3.3, 10.1,
         "A useful home / lab robot has to take a vague, natural-language request and "
         "carry it out end-to-end — understand it, explore, find the object, and physically fetch it.",
         fs=20, bold=False, color=MUTE, align=PP_ALIGN.CENTER)
placeholder(s, 3.5, 4.7, 6.3, 2.0, "[ scene photo: robot in the room ]")
footer(s, 2)

# ============================================================ SLIDE 3 — Why it's hard
s = slide()
title(s, "Why it’s hard")
bullets(s, 0.9, 1.9, 11.5, [
    ("Long-horizon", "“go to A, then B” — chained navigation + manipulation, not one step"),
    ("Partially observable", "objects and obstacles are unknown until the robot actually sees them"),
    ("Real hardware, not simulation", "noisy depth, reflective / transparent objects, limited arm reach"),
    ("One closed loop", "language → navigation → perception → grasp → memory, all tied together"),
], fs=23, gap=8)
footer(s, 3)

# ============================================================ SLIDE 4 — Platform
s = slide()
title(s, "The Platform")
bullets(s, 0.9, 1.9, 6.0, [
    ("Mobile base", "4-wheel mecanum — omnidirectional driving"),
    ("Arm", "5-DOF manipulator + parallel gripper"),
    ("Sensing", "Orbbec DaBai DCW RGB-D depth camera (eye-in-hand)"),
    ("Compute", "onboard Jetson, ROS control stack"),
], fs=20, gap=8)
placeholder(s, 7.4, 1.9, 5.3, 4.5, "[ robot photo ]")
footer(s, 4)

# ============================================================ SLIDE 5 — Architecture (money slide)
s = slide()
title(s, "System Architecture")
# top inputs
box(s, 2.95, 1.10, 2.5, 0.62, "Instruction", C_BLUE, fs=15, bold=False)
box(s, 8.20, 1.10, 2.6, 0.62, "Observations\n(RGB-D + odometry)", C_PURPLE, fs=12, bold=False)
# center column
box(s, 2.95, 2.05, 2.5, 0.68, "LLM Planner", C_GRAY, fs=16, bold=True)
box(s, 2.95, 3.02, 2.5, 0.58, "Subtask Plan", C_GREEN, fs=15, bold=True)
box(s, 2.95, 3.90, 2.5, 0.66, "Skill Executor", C_GRAY, fs=15, bold=True)
# skills
box(s, 1.55, 4.95, 1.95, 0.62, "Navigation\n(SLAM + Nav2)", C_YELLOW, fs=12)
box(s, 4.90, 4.95, 1.95, 0.62, "Depth Grasping\n(hand-eye + IK)", C_PINK, fs=12)
# controller + env
box(s, 2.55, 5.92, 3.3, 0.55, "Controller  /  ROS 2", C_TAN, fs=14)
box(s, 2.55, 6.62, 3.3, 0.50, "Robot  +  Sensors   (Environment)", C_PINK, fs=13)
# right feedback column
box(s, 8.30, 3.90, 2.5, 0.66, "Perception\n(depth → object pose)", C_ORANGE, fs=12)
box(s, 8.30, 2.90, 2.5, 0.62, "Memory: bottles.yaml", C_GREEN, fs=13, bold=True)
# arrows (center chain)
arrow(s, 4.20, 1.72, 4.20, 2.05)          # instruction -> planner
arrow(s, 4.20, 2.73, 4.20, 3.02)          # planner -> subtasks
arrow(s, 4.20, 3.60, 4.20, 3.90)          # subtasks -> executor
arrow(s, 3.70, 4.56, 2.53, 4.95)          # executor -> nav
arrow(s, 4.70, 4.56, 5.87, 4.95)          # executor -> grasp
arrow(s, 2.53, 5.57, 3.30, 5.92)          # nav -> controller
arrow(s, 5.87, 5.57, 5.10, 5.92)          # grasp -> controller
arrow(s, 4.20, 6.47, 4.20, 6.62)          # controller -> robot
# feedback path: robot -> perception -> memory -> planner ; observations -> planner
arrow(s, 5.85, 6.87, 9.55, 6.87); arrow(s, 9.55, 6.87, 9.55, 4.56)  # robot obs up
arrow(s, 9.55, 3.90, 9.55, 3.52)          # perception -> memory
arrow(s, 8.30, 3.21, 5.45, 2.55)          # memory -> planner (re-plan)
arrow(s, 8.20, 1.41, 5.45, 2.20)          # observations -> planner
alabel(s, 4.05, 2.75, "plan")
alabel(s, 6.6, 6.55, "obs")
alabel(s, 6.5, 2.55, "re-plan")
# side notes
note(s, 0.35, 2.00, 2.35, 1.15,
     "LLM decomposes the instruction into ordered subtasks and re-plans as new objects are seen")
note(s, 10.95, 3.75, 2.05, 1.15,
     "Depth + detection turn raw RGB-D into 3-D object poses")
note(s, 10.95, 2.55, 2.05, 1.05,
     "Detections persist to bottles.yaml → long-horizon, resumable tasks")
footer(s, 5)

# ============================================================ SLIDE 6 — Navigation
s = slide()
title(s, "Navigation: long-horizon & map-aware")
bullets(s, 0.9, 1.9, 6.1, [
    ("SLAM + Nav2", "Hector SLAM builds the map; Nav2 plans and drives to goals"),
    ("Multi-hop goals", "“go to A, then B” verified on hardware — two hops, both reached"),
    ("Reactive to the world", "obstacles that enter the map are planned around"),
    ("Tuned for the real bot", "throttled teleop + path-bias tuning to stay stable, avoid map corruption"),
], fs=19, gap=8)
placeholder(s, 7.4, 1.9, 5.3, 4.5, "[ rviz map / trajectory screenshot ]")
footer(s, 6)

# ============================================================ SLIDE 7 — Grasping
s = slide()
title(s, "Depth-based Grasping")
bullets(s, 0.9, 1.75, 6.2, [
    ("Pipeline", "RGB-D → object detection → hand-eye (eye-in-hand) → IK → approach / descend / close / lift"),
    ("First pick on hardware", "lifted a wrapped can off the table end-to-end"),
    ("Hard problems solved", "in-process FK/IK when the vendor IK service returned no solution; "
                              "nav + grasp share one control node without contending for the serial bus"),
    ("Open challenge", "reflective / transparent bottles lose depth — robust depth sampling in progress"),
], fs=17, gap=7)
placeholder(s, 7.6, 1.9, 5.1, 4.5, "[ grasp photo / depth view ]")
footer(s, 7)

# ============================================================ SLIDE 8 — Memory / closed loop
s = slide()
title(s, "Memory closes the loop")
bullets(s, 0.9, 2.0, 11.5, [
    ("Persist what you see", "every detected object is written to bottles.yaml with its pose"),
    ("Plan over memory", "the LLM filters candidates and orders multi-point navigation to collect them"),
    ("Why it matters", "this is what turns one-shot reactions into long-horizon, resumable tasks"),
], fs=23, gap=12)
placeholder(s, 3.4, 4.9, 6.5, 1.7, "[ bottles.yaml → planned route diagram ]")
footer(s, 8)

# ============================================================ SLIDE 9 — Experiments
s = slide()
title(s, "Experiments — 5-bottle task")
box(s, 0.9, 2.0, 5.6, 1.9,
    "Experiment A  —  targeted collection\n\nTask: fetch all the Coke bottles\n\nResult:  3 / 3 collected",
    C_GREEN, fs=17, align=PP_ALIGN.LEFT)
box(s, 6.85, 2.0, 5.6, 1.9,
    "Experiment B  —  full patrol\n\nTask: visit and log every bottle\n\nResult:  5 / 5 found",
    C_BLUE, fs=17, align=PP_ALIGN.LEFT)
box(s, 0.9, 4.3, 11.55, 1.7,
    "Milestone (2026-07-08):  first full  perception → planning → action → memory  loop "
    "running on the real robot — LLM candidate filtering + multi-point navigation + on-disk memory.",
    C_YELLOW, fs=18, align=PP_ALIGN.LEFT)
footer(s, 9)

# ============================================================ SLIDE 10 — Demo
s = slide()
title(s, "Demo")
placeholder(s, 0.9, 2.0, 5.6, 4.2, "[ video / GIF ]\nExp A — collect Coke")
placeholder(s, 6.85, 2.0, 5.6, 4.2, "[ video / GIF ]\nExp B — patrol")
footer(s, 10)

# ============================================================ SLIDE 11 — Takeaways
s = slide()
title(s, "Key Takeaways")
bullets(s, 0.9, 1.85, 11.6, [
    ("Full-stack robotics, on real hardware",
     "ROS, Hector SLAM + Nav2, hand-eye calibration, LLM planning + memory, systems debugging"),
    ("Hardest wins",
     "nav + grasp coexistence on a single control node; in-process IK; robust depth on tough surfaces"),
    ("Honest next steps",
     "reliable grasp on reflective bottles; tighter reach calibration"),
], fs=21, gap=10)
contact = s.shapes.add_textbox(Inches(0.9), Inches(6.15), Inches(11.5), Inches(0.6))
p = contact.text_frame.paragraphs[0]
r = p.add_run(); r.text = "[ your email ]   ·   [ github / demo link ]"
r.font.name = FONT; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = BAR
footer(s, 11)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "JetRover_teaser.pptx")
prs.save(out)
print("saved:", out, "  slides:", len(prs.slides._sldIdLst))
